#!/usr/bin/env python3

from __future__ import annotations

import argparse
import csv
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SCENARIO_LABELS = {
    "updates_100": "100% updates",
    "balanced_50_50": "50% updates / 50% scans",
    "scans_90": "10% updates / 90% scans",
    "scans_100": "100% scans",
}

IMPL_LABELS = {
    "snapshot_baseline": "Baseline snapshot",
    "mcas_volatile": "MCAS snapshot (volatile)",
}

IMPL_COLORS = {
    "snapshot_baseline": RGBColor(56, 189, 248),
    "mcas_volatile": RGBColor(167, 139, 250),
}


def repo_root() -> Path:
    return Path(__file__).resolve().parents[2]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the benchmark presentation deck.")
    parser.add_argument(
        "--input-csv",
        default=str(repo_root() / "presentation" / "data" / "benchmark_results_latest.csv"),
        help="CSV file produced by the benchmark parser.",
    )
    parser.add_argument(
        "--assets-dir",
        default=str(repo_root() / "presentation" / "assets"),
        help="Directory containing generated chart images.",
    )
    parser.add_argument(
        "--output",
        default=str(repo_root() / "presentation" / "mcas_snapshot_benchmark_presentation.pptx"),
        help="Output .pptx file.",
    )
    return parser.parse_args()


def read_rows(path: Path) -> list[dict[str, object]]:
    rows: list[dict[str, object]] = []
    with path.open("r", newline="", encoding="utf-8") as handle:
        reader = csv.DictReader(handle)
        for row in reader:
            rows.append(
                {
                    "impl_name": row["impl_name"],
                    "scenario_name": row["scenario_name"],
                    "domains": int(row["domains"]),
                    "total_ops": int(row["total_ops"]),
                    "seconds": float(row["seconds"]),
                    "avg_us_per_op": float(row["avg_us_per_op"]),
                }
            )
    return rows


def add_background(slide) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(15, 23, 42)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(30, 41, 59)
    band.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(9.8), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(248, 250, 252)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.1), Inches(10.5), Inches(0.45))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(148, 163, 184)


def add_bullet_list(slide, items: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(10)


def add_caption(slide, text: str, left: float, top: float, width: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(148, 163, 184)


def add_metric_card(slide, title: str, value: str, left: float, top: float, width: float) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(1.15),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.color.rgb = RGBColor(51, 65, 85)

    title_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.3), Inches(0.3))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(148, 163, 184)

    value_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.42), Inches(width - 0.3), Inches(0.4))
    p = value_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor(248, 250, 252)


def human_ms(value: float) -> str:
    return f"{value * 1000:.1f} ms"


def human_us(value: float) -> str:
    return f"{value:.2f} us/op"


def group_rows(rows: list[dict[str, object]], min_domains: int = 2) -> dict[str, dict[str, list[dict[str, object]]]]:
    grouped: dict[str, dict[str, list[dict[str, object]]]] = defaultdict(lambda: defaultdict(list))
    for row in rows:
        if int(row["domains"]) < min_domains:
            continue
        grouped[str(row["scenario_name"])][str(row["impl_name"])].append(row)
    for scenario_rows in grouped.values():
        for impl_rows in scenario_rows.values():
            impl_rows.sort(key=lambda item: int(item["domains"]))
    return grouped


def top_findings(grouped: dict[str, dict[str, list[dict[str, object]]]]) -> list[str]:
    balanced = grouped["balanced_50_50"]
    baseline = float(balanced["snapshot_baseline"][-1]["seconds"])
    volatile = float(balanced["mcas_volatile"][-1]["seconds"])

    scans = grouped["scans_100"]
    scan_baseline = float(scans["snapshot_baseline"][-1]["seconds"])
    scan_volatile = float(scans["mcas_volatile"][-1]["seconds"])

    updates = grouped["updates_100"]
    update_baseline = float(updates["snapshot_baseline"][-1]["seconds"])
    update_volatile = float(updates["mcas_volatile"][-1]["seconds"])

    return [
        f"At 8 threads on a balanced workload, the baseline finishes in {human_ms(baseline)} while the MCAS snapshot takes {human_ms(volatile)}.",
        f"For 100% scans at 8 threads, the MCAS snapshot is {scan_volatile / scan_baseline:.1f}x slower than the double-collect baseline.",
        f"Even on 100% updates at 8 threads, the MCAS layer still stretches total runtime from {human_ms(update_baseline)} to {human_ms(update_volatile)}.",
    ]


def scenario_findings(grouped: dict[str, dict[str, list[dict[str, object]]]], scenario: str) -> list[str]:
    scenario_rows = grouped[scenario]
    baseline = float(scenario_rows["snapshot_baseline"][-1]["seconds"])
    volatile = float(scenario_rows["mcas_volatile"][-1]["seconds"])
    baseline_avg = float(scenario_rows["snapshot_baseline"][-1]["avg_us_per_op"])
    volatile_avg = float(scenario_rows["mcas_volatile"][-1]["avg_us_per_op"])
    return [
        f"At 8 threads, baseline elapsed time is {human_ms(baseline)}.",
        f"At 8 threads, MCAS elapsed time is {human_ms(volatile)} ({volatile / baseline:.1f}x baseline time).",
        f"Average cost per operation grows from {human_us(baseline_avg)} to {human_us(volatile_avg)}.",
    ]


def add_image_slide(slide, title: str, subtitle: str, image_path: Path, bullets: list[str]) -> None:
    add_background(slide)
    add_title(slide, title, subtitle)
    slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.7), width=Inches(8.2))
    add_bullet_list(slide, bullets, left=9.2, top=2.0, width=3.3, height=3.8)
    add_caption(slide, "Plots use the required 2-8 thread range from the assignment deliverable.", 0.7, 6.95, 7.2)


def build_deck(
    rows: list[dict[str, object]],
    assets_dir: Path,
    output_path: Path,
    input_csv_name: str,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    grouped = group_rows(rows, min_domains=2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        "Multi-Word Compare-and-Swap in OCaml 5",
        "Benchmarking and snapshot-performance presentation deck",
    )
    hero_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(6.8), Inches(2.2))
    p = hero_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "How much simplicity does MCAS buy us, and what does the software layer cost?"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(248, 250, 252)
    add_bullet_list(
        slide,
        [
            "Implements descriptor-based MCAS and applies it to atomic snapshot.",
            "Benchmarks baseline snapshot against the volatile MCAS snapshot only.",
            "Slides focus on elapsed time trends across the benchmark deliverable.",
        ],
        left=0.8,
        top=4.0,
        width=6.7,
        height=2.1,
    )
    add_metric_card(slide, "Workloads", "4 ratios", 8.8, 2.1, 1.7)
    add_metric_card(slide, "Thread range", "2-8", 10.7, 2.1, 1.7)
    add_metric_card(slide, "Implementations", "2", 8.8, 3.55, 1.7)
    add_metric_card(slide, "Registers", "16", 10.7, 3.55, 1.7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Deliverables and Method", "What this benchmark section is measuring")
    add_bullet_list(
        slide,
        [
            "Deliverable focus: compare snapshot elapsed time across varying update/scan ratios and 2-8 threads.",
            "Workloads: 100% updates, 50/50, 10/90, and 100% scans.",
            "Implementations: baseline double-collect snapshot and volatile MCAS snapshot.",
            "Benchmark harness uses 16 registers and 2,000 operations per domain.",
            "Each point reports total elapsed time for the full workload, plus average microseconds per operation.",
        ],
        left=0.8,
        top=1.8,
        width=11.6,
        height=4.6,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Benchmark Overview", "All workload mixes in one view")
    slide.shapes.add_picture(str(assets_dir / "benchmark_overview.png"), Inches(0.7), Inches(1.55), width=Inches(8.7))
    slide.shapes.add_picture(str(assets_dir / "relative_to_baseline.png"), Inches(9.55), Inches(1.9), width=Inches(3.0))
    add_bullet_list(slide, top_findings(grouped), left=0.8, top=6.0, width=12.0, height=1.0)

    for scenario in SCENARIO_LABELS:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_image_slide(
            slide,
            SCENARIO_LABELS[scenario],
            "Elapsed time across the required thread counts",
            assets_dir / f"{scenario}.png",
            scenario_findings(grouped, scenario),
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Takeaways", "What the benchmark results say about the project")
    add_bullet_list(
        slide,
        [
            "MCAS makes the snapshot implementation conceptually uniform: update and scan are both expressed in terms of MCAS operations.",
            "That expressiveness comes with a visible time cost, especially for scan-heavy workloads where double-collect finishes much faster.",
            "The time-based view is easier to explain in a presentation: same workload mix, more wall-clock time once the MCAS software layer is involved.",
            "The cleanest story for the deck is simplicity versus runtime cost: MCAS gives a neat abstraction, but not a free one.",
        ],
        left=0.8,
        top=1.8,
        width=11.7,
        height=4.4,
    )
    add_caption(slide, f"Deck generated from parsed benchmark data in {input_csv_name}.", 0.8, 6.95, 9.0)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> int:
    args = parse_args()
    rows = read_rows(Path(args.input_csv))
    build_deck(rows, Path(args.assets_dir), Path(args.output), Path(args.input_csv).name)
    print(f"Generated presentation deck at {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
